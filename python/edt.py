import sys
from pptx import Presentation
import shutil

input_ppt = sys.argv[1]
output_ppt = sys.argv[2]
name = sys.argv[3]

# Copy template → output
shutil.copyfile(input_ppt, output_ppt)

prs = Presentation(output_ppt)

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = shape.text.replace("{{NAME}}", name)

prs.save(output_ppt)

print("PPT updated successfully")
