from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import sys
import os
import shutil

# -------- Arguments from Node --------
template_ppt = sys.argv[1]   # uploads/input.pptx (READ ONLY)
output_ppt = sys.argv[2]     # uploads/ppt/<id>/kickoff_<id>.pptx
name_to_add = sys.argv[3]    # Customer name

print("TEMPLATE:", template_ppt)
print("OUTPUT:", output_ppt)
print("NAME:", name_to_add)

if not os.path.exists(template_ppt):
    print("Template file not found")
    sys.exit(1)

# ✅ STEP 1: COPY TEMPLATE → OUTPUT
shutil.copyfile(template_ppt, output_ppt)
print("Template copied to output path")

# ✅ STEP 2: OPEN COPIED FILE (NOT TEMPLATE)
prs = Presentation(output_ppt)
slide = prs.slides[0]

# Collect existing text to avoid duplicates
existing_text = []
for shape in slide.shapes:
    if shape.has_text_frame:
        existing_text.append(shape.text.strip())

if name_to_add not in existing_text:
    # Slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Textbox size
    textbox_width = Inches(4)
    textbox_height = Inches(0.5)

    # Bottom-left position
    left = Inches(1.5)
    top = slide_height - textbox_height - Inches(4.0)

    textbox = slide.shapes.add_textbox(
        left, top, textbox_width, textbox_height
    )
    textbox.text = name_to_add

    # Font styling
    for paragraph in textbox.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(255, 255, 255)

    print(f"Added name at bottom-left in white: {name_to_add}")
else:
    print(f"Name already exists: {name_to_add}")

# ✅ STEP 3: SAVE ONLY OUTPUT FILE
prs.save(output_ppt)
print(f"Saved updated file: {output_ppt}")
